"""PPT node — download → slide text dump → Claude."""

from __future__ import annotations

import asyncio
import logging
from pathlib import Path

from core.downloader import download_to_temp
from core.llm import extract_from_content
from models import ExtractionError, ExtractedContent, FileItem, FileType
from state import CIMState

log = logging.getLogger(__name__)


def _to_text(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


async def _process_one(item: FileItem, listing_xml: str = "") -> ExtractedContent:
    try:
        if item.content:
            import tempfile, base64 as _b64
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(_b64.b64decode(item.content))
                tmp_path = Path(tmp.name)
            try:
                raw_text = await asyncio.to_thread(_to_text, tmp_path)
            finally:
                tmp_path.unlink(missing_ok=True)
        else:
            async with download_to_temp(item.url) as path:
                raw_text = await asyncio.to_thread(_to_text, path)

        content_blocks = [{"type": "text", "text": raw_text}]
        text = await extract_from_content(content_blocks, item.url, listing_xml)
        log.info("  PPT extracted: %r", item.label)
        return ExtractedContent(
            source_url=item.url,
            file_type=FileType.PPT,
            metadata={"extracted_text": text},
        )
    except Exception as exc:
        log.error("  PPT failed: %r — %s", item.label, exc)
        return ExtractedContent(source_url=item.url, file_type=FileType.PPT, error=str(exc))


async def ppt_node(state: CIMState) -> dict:
    files: list[FileItem] = state.get("ppt_files", [])
    if not files:
        return {"extracted": [], "errors": []}

    log.info("PPT: processing %d file(s)", len(files))
    results = await asyncio.gather(*[_process_one(f, state.get("listing_xml", "")) for f in files])

    extracted = [r for r in results if not r.error]
    errors = [ExtractionError(url=r.source_url, stage="ppt", message=r.error) for r in results if r.error]
    if errors:
        log.error("PPT: %d error(s)", len(errors))
    return {"extracted": list(extracted), "errors": errors}
